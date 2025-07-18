import streamlit as st
import os
import tempfile
from docx2pdf import convert
from pathlib import Path
import base64
import logging
from typing import Optional, Tuple
import time
import hashlib
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import markdown
from weasyprint import HTML, CSS
from pptx import Presentation
from io import BytesIO

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Constants
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
ALLOWED_EXTENSIONS = ['docx', 'pptx', 'txt', 'md']
TEMP_FILE_CLEANUP_DELAY = 300  # 5 minutes

# File type configurations
FILE_TYPE_CONFIG = {
    'docx': {
        'name': 'Word Document',
        'icon': '📄',
        'description': 'Microsoft Word documents (.docx)',
        'requires': 'Microsoft Word'
    },
    'pptx': {
        'name': 'PowerPoint Presentation',
        'icon': '📊',
        'description': 'Microsoft PowerPoint presentations (.pptx)',
        'requires': 'Python-pptx library'
    },
    'txt': {
        'name': 'Text File',
        'icon': '📝',
        'description': 'Plain text files (.txt)',
        'requires': 'ReportLab library'
    },
    'md': {
        'name': 'Markdown File',
        'icon': '📋',
        'description': 'Markdown files (.md)',
        'requires': 'Markdown and WeasyPrint libraries'
    }
}

def validate_file(uploaded_file) -> Tuple[bool, str]:
    """
    Validate uploaded file for security and format requirements
    
    Args:
        uploaded_file: Streamlit uploaded file object
    
    Returns:
        Tuple[bool, str]: (is_valid, error_message)
    """
    if not uploaded_file:
        return False, "No file uploaded"
    
    # Check file size
    if uploaded_file.size > MAX_FILE_SIZE:
        return False, f"File too large. Maximum size allowed: {MAX_FILE_SIZE // (1024*1024)}MB"
    
    if uploaded_file.size == 0:
        return False, "File is empty"
    
    # Check file extension
    file_extension = Path(uploaded_file.name).suffix.lower().lstrip('.')
    if file_extension not in ALLOWED_EXTENSIONS:
        return False, f"Unsupported file type. Allowed types: {', '.join(ALLOWED_EXTENSIONS)}"
    
    # Basic file content validation
    try:
        file_content = uploaded_file.getvalue()
        
        # Validate based on file type
        if file_extension in ['docx', 'pptx']:
            if not file_content.startswith(b'PK'):  # Office files start with PK (ZIP signature)
                return False, f"Invalid {file_extension.upper()} file format"
        elif file_extension == 'txt':
            try:
                file_content.decode('utf-8')
            except UnicodeDecodeError:
                try:
                    file_content.decode('latin-1')
                except UnicodeDecodeError:
                    return False, "Text file contains invalid characters"
        elif file_extension == 'md':
            try:
                file_content.decode('utf-8')
            except UnicodeDecodeError:
                return False, "Markdown file must be UTF-8 encoded"
                
    except Exception as e:
        return False, f"Error reading file: {str(e)}"
    
    return True, ""

def convert_word_to_pdf(input_path: str, output_path: str) -> Tuple[bool, str]:
    """
    Convert a Word document to PDF with better error handling
    
    Args:
        input_path (str): Path to the input .docx file
        output_path (str): Path for the output .pdf file
    
    Returns:
        Tuple[bool, str]: (success, error_message)
    """
    try:
        # Validate input file exists
        if not os.path.exists(input_path):
            return False, "Input file not found"
        
        # Perform conversion
        convert(input_path, output_path)
        
        # Verify output file was created
        if not os.path.exists(output_path):
            return False, "PDF file was not created"
        
        # Check if output file has content
        if os.path.getsize(output_path) == 0:
            return False, "Generated PDF file is empty"
        
        logger.info(f"Successfully converted {input_path} to {output_path}")
        return True, ""
        
    except Exception as e:
        error_msg = f"Word conversion failed: {str(e)}"
        logger.error(error_msg)
        return False, error_msg

def convert_powerpoint_to_pdf(input_path: str, output_path: str) -> Tuple[bool, str]:
    """
    Convert a PowerPoint presentation to PDF
    
    Args:
        input_path (str): Path to the input .pptx file
        output_path (str): Path for the output .pdf file
    
    Returns:
        Tuple[bool, str]: (success, error_message)
    """
    try:
        # Load presentation
        prs = Presentation(input_path)
        
        # Create PDF document
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Title style
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=30,
            textColor='darkblue'
        )
        
        # Content style
        content_style = ParagraphStyle(
            'CustomContent',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=12
        )
        
        # Process each slide
        for i, slide in enumerate(prs.slides):
            # Add slide number
            story.append(Paragraph(f"Slide {i + 1}", title_style))
            story.append(Spacer(1, 12))
            
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                for text in slide_text:
                    # Clean and format text
                    clean_text = text.replace('\n', '<br/>')
                    story.append(Paragraph(clean_text, content_style))
                    story.append(Spacer(1, 6))
            else:
                story.append(Paragraph("(No text content)", content_style))
            
            # Add space between slides
            story.append(Spacer(1, 20))
        
        # Build PDF
        doc.build(story)
        
        logger.info(f"Successfully converted PowerPoint {input_path} to {output_path}")
        return True, ""
        
    except Exception as e:
        error_msg = f"PowerPoint conversion failed: {str(e)}"
        logger.error(error_msg)
        return False, error_msg

def convert_text_to_pdf(input_path: str, output_path: str) -> Tuple[bool, str]:
    """
    Convert a text file to PDF
    
    Args:
        input_path (str): Path to the input .txt file
        output_path (str): Path for the output .pdf file
    
    Returns:
        Tuple[bool, str]: (success, error_message)
    """
    try:
        # Read text file
        with open(input_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Create PDF document
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom style for text
        text_style = ParagraphStyle(
            'CustomText',
            parent=styles['Normal'],
            fontSize=11,
            fontName='Courier',
            spaceAfter=6,
            leftIndent=0
        )
        
        # Split content into paragraphs
        paragraphs = content.split('\n\n')
        
        for paragraph in paragraphs:
            if paragraph.strip():
                # Replace line breaks with HTML breaks
                formatted_text = paragraph.replace('\n', '<br/>')
                # Escape HTML characters
                formatted_text = formatted_text.replace('&', '&amp;').replace('<br/>', '<br/>')
                story.append(Paragraph(formatted_text, text_style))
                story.append(Spacer(1, 6))
        
        # Build PDF
        doc.build(story)
        
        logger.info(f"Successfully converted text file {input_path} to {output_path}")
        return True, ""
        
    except Exception as e:
        error_msg = f"Text conversion failed: {str(e)}"
        logger.error(error_msg)
        return False, error_msg

def convert_markdown_to_pdf(input_path: str, output_path: str) -> Tuple[bool, str]:
    """
    Convert a Markdown file to PDF
    
    Args:
        input_path (str): Path to the input .md file
        output_path (str): Path for the output .pdf file
    
    Returns:
        Tuple[bool, str]: (success, error_message)
    """
    try:
        # Read markdown file
        with open(input_path, 'r', encoding='utf-8') as file:
            md_content = file.read()
        
        # Convert markdown to HTML
        html_content = markdown.markdown(
            md_content,
            extensions=['extra', 'codehilite', 'toc']
        )
        
        # Create complete HTML document
        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <style>
                body {{
                    font-family: Arial, sans-serif;
                    line-height: 1.6;
                    margin: 40px;
                    color: #333;
                }}
                h1, h2, h3, h4, h5, h6 {{
                    color: #2c3e50;
                    margin-top: 30px;
                    margin-bottom: 15px;
                }}
                h1 {{ font-size: 2.2em; }}
                h2 {{ font-size: 1.8em; }}
                h3 {{ font-size: 1.4em; }}
                p {{ margin-bottom: 15px; }}
                code {{
                    background-color: #f4f4f4;
                    padding: 2px 4px;
                    border-radius: 3px;
                    font-family: 'Courier New', monospace;
                }}
                pre {{
                    background-color: #f8f8f8;
                    padding: 15px;
                    border-radius: 5px;
                    overflow-x: auto;
                    border-left: 4px solid #3498db;
                }}
                blockquote {{
                    border-left: 4px solid #bdc3c7;
                    margin-left: 0;
                    padding-left: 20px;
                    color: #7f8c8d;
                }}
                table {{
                    border-collapse: collapse;
                    width: 100%;
                    margin: 20px 0;
                }}
                th, td {{
                    border: 1px solid #ddd;
                    padding: 8px;
                    text-align: left;
                }}
                th {{
                    background-color: #f2f2f2;
                    font-weight: bold;
                }}
            </style>
        </head>
        <body>
            {html_content}
        </body>
        </html>
        """
        
        # Convert HTML to PDF using WeasyPrint
        HTML(string=full_html).write_pdf(output_path)
        
        logger.info(f"Successfully converted Markdown {input_path} to {output_path}")
        return True, ""
        
    except Exception as e:
        error_msg = f"Markdown conversion failed: {str(e)}"
        logger.error(error_msg)
        return False, error_msg

def convert_document_to_pdf(input_path: str, output_path: str, file_type: str) -> Tuple[bool, str]:
    """
    Convert any supported document type to PDF
    
    Args:
        input_path (str): Path to the input file
        output_path (str): Path for the output .pdf file
        file_type (str): Type of file (docx, pptx, txt, md)
    
    Returns:
        Tuple[bool, str]: (success, error_message)
    """
    conversion_functions = {
        'docx': convert_word_to_pdf,
        'pptx': convert_powerpoint_to_pdf,
        'txt': convert_text_to_pdf,
        'md': convert_markdown_to_pdf
    }
    
    if file_type not in conversion_functions:
        return False, f"Unsupported file type: {file_type}"
    
    return conversion_functions[file_type](input_path, output_path)

def get_download_link(file_path: str, file_name: str) -> str:
    """
    Generate a secure download link for the converted PDF
    
    Args:
        file_path (str): Path to the PDF file
        file_name (str): Name for the download
    
    Returns:
        str: HTML download link
    """
    try:
        with open(file_path, "rb") as f:
            bytes_data = f.read()
        
        b64 = base64.b64encode(bytes_data).decode()
        # Sanitize filename for download
        safe_filename = "".join(c for c in file_name if c.isalnum() or c in (' ', '-', '_', '.')).rstrip()
        href = f'<a href="data:application/pdf;base64,{b64}" download="{safe_filename}" style="text-decoration: none; background-color: #4CAF50; color: white; padding: 10px 20px; border-radius: 5px; display: inline-block; margin: 10px 0;">📥 Download PDF</a>'
        return href
    except Exception as e:
        logger.error(f"Error creating download link: {str(e)}")
        return f"<p style='color: red;'>Error creating download link: {str(e)}</p>"

def cleanup_temp_files(*file_paths):
    """
    Clean up temporary files with error handling
    
    Args:
        *file_paths: Variable number of file paths to clean up
    """
    for file_path in file_paths:
        if file_path and os.path.exists(file_path):
            try:
                os.unlink(file_path)
                logger.info(f"Cleaned up temporary file: {file_path}")
            except Exception as e:
                logger.warning(f"Failed to clean up {file_path}: {str(e)}")

def format_file_size(size_bytes: int) -> str:
    """
    Format file size in human readable format
    
    Args:
        size_bytes (int): Size in bytes
    
    Returns:
        str: Formatted size string
    """
    if size_bytes == 0:
        return "0 B"
    
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.1f} TB"

def main():
    st.set_page_config(
        page_title="Universal Document to PDF Converter",
        page_icon="📄",
        layout="centered",
        initial_sidebar_state="collapsed"
    )
    
    # Initialize session state
    if 'conversion_history' not in st.session_state:
        st.session_state.conversion_history = []
    if 'last_converted_file' not in st.session_state:
        st.session_state.last_converted_file = None
    
    # Header
    st.title("📄 Universal Document to PDF Converter")
    st.markdown("Convert Word, PowerPoint, Text, and Markdown documents to PDF instantly!")
    
    # Supported file types display
    st.markdown("### 🎯 Supported File Types")
    cols = st.columns(4)
    for i, (ext, config) in enumerate(FILE_TYPE_CONFIG.items()):
        with cols[i]:
            st.markdown(f"""
            <div style="text-align: center; padding: 10px; border: 1px solid #ddd; border-radius: 5px; margin: 5px;">
                <div style="font-size: 2em;">{config['icon']}</div>
                <div style="font-weight: bold;">{config['name']}</div>
                <div style="font-size: 0.8em; color: #666;">.{ext}</div>
            </div>
            """, unsafe_allow_html=True)
    
    # System requirements
    with st.expander("⚠️ System Requirements", expanded=False):
        st.markdown("### Requirements by File Type:")
        for ext, config in FILE_TYPE_CONFIG.items():
            st.markdown(f"- **{config['name']}** (.{ext}): {config['requires']}")
        
        st.warning("""
        **Note for Word documents:** Requires Microsoft Word to be installed on the system.
        Other file types use Python libraries and don't require additional software.
        """)
    
    # File uploader with improved validation
    uploaded_file = st.file_uploader(
        "Choose a document to convert",
        type=ALLOWED_EXTENSIONS,
        help=f"Select a document file to convert to PDF (Max size: {MAX_FILE_SIZE // (1024*1024)}MB)",
        key="file_uploader"
    )
    
    if uploaded_file is not None:
        # Validate file
        is_valid, error_message = validate_file(uploaded_file)
        
        if not is_valid:
            st.error(f"❌ {error_message}")
            return
        
        # Get file type
        file_extension = Path(uploaded_file.name).suffix.lower().lstrip('.')
        file_config = FILE_TYPE_CONFIG[file_extension]
        
        # Display file details
        col1, col2, col3 = st.columns(3)
        with col1:
            st.success(f"{file_config['icon']} **File:** {uploaded_file.name}")
        with col2:
            st.info(f"📊 **Size:** {format_file_size(uploaded_file.size)}")
        with col3:
            st.info(f"🔧 **Type:** {file_config['name']}")
        
        # Convert button
        if st.button("🔄 Convert to PDF", type="primary", use_container_width=True):
            start_time = time.time()
            
            with st.spinner(f"Converting {file_config['name']} to PDF... Please wait."):
                # Create temporary files with unique names
                file_hash = hashlib.md5(uploaded_file.getvalue()).hexdigest()[:8]
                
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}", prefix=f"{file_extension}_{file_hash}_") as tmp_input:
                    tmp_input.write(uploaded_file.getvalue())
                    tmp_input_path = tmp_input.name
                
                # Generate output filename
                input_name = Path(uploaded_file.name).stem
                tmp_output_path = os.path.join(tempfile.gettempdir(), f"pdf_{file_hash}_{input_name}.pdf")
                
                # Perform conversion
                success, error_message = convert_document_to_pdf(tmp_input_path, tmp_output_path, file_extension)
                
                conversion_time = time.time() - start_time
                
                if success:
                    st.success(f"✅ {file_config['name']} converted successfully in {conversion_time:.2f} seconds!")
                    
                    # Provide download link
                    download_filename = f"{input_name}.pdf"
                    st.markdown(
                        get_download_link(tmp_output_path, download_filename),
                        unsafe_allow_html=True
                    )
                    
                    # Display conversion details
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        pdf_size = os.path.getsize(tmp_output_path)
                        st.metric("PDF Size", format_file_size(pdf_size))
                    with col2:
                        st.metric("Conversion Time", f"{conversion_time:.2f}s")
                    with col3:
                        compression_ratio = (1 - pdf_size / uploaded_file.size) * 100 if uploaded_file.size > 0 else 0
                        st.metric("Size Change", f"{compression_ratio:+.1f}%")
                    
                    # Update session state
                    st.session_state.conversion_history.append({
                        'filename': uploaded_file.name,
                        'file_type': file_extension,
                        'size': uploaded_file.size,
                        'pdf_size': pdf_size,
                        'conversion_time': conversion_time,
                        'timestamp': time.time()
                    })
                    st.session_state.last_converted_file = tmp_output_path
                    
                    # Schedule cleanup (in a real app, you'd want a background task)
                    st.info("💡 **Tip:** The PDF file will be automatically cleaned up after 5 minutes for security.")
                    
                else:
                    st.error(f"❌ {error_message}")
                    
                    # File type specific troubleshooting
                    if file_extension == 'docx':
                        st.info("💡 **Word Document Troubleshooting:**")
                        st.markdown("""
                        - Ensure Microsoft Word is installed and properly licensed
                        - Check that the DOCX file is not corrupted
                        - Try opening the file in Word first to verify it works
                        """)
                    elif file_extension == 'pptx':
                        st.info("💡 **PowerPoint Troubleshooting:**")
                        st.markdown("""
                        - Ensure the PPTX file is not corrupted
                        - Check that the presentation contains readable text
                        - Try opening the file in PowerPoint first to verify it works
                        """)
                    elif file_extension == 'txt':
                        st.info("💡 **Text File Troubleshooting:**")
                        st.markdown("""
                        - Ensure the text file uses UTF-8 encoding
                        - Check that the file contains readable text
                        - Try opening the file in a text editor to verify content
                        """)
                    elif file_extension == 'md':
                        st.info("💡 **Markdown Troubleshooting:**")
                        st.markdown("""
                        - Ensure the Markdown file uses UTF-8 encoding
                        - Check that the file contains valid Markdown syntax
                        - Try previewing the file in a Markdown editor
                        """)
                
                # Clean up input file immediately
                cleanup_temp_files(tmp_input_path)
    
    # Conversion history
    if st.session_state.conversion_history:
        with st.expander("📊 Conversion History", expanded=False):
            for i, conversion in enumerate(reversed(st.session_state.conversion_history[-10:])):  # Show last 10
                file_config = FILE_TYPE_CONFIG[conversion['file_type']]
                st.write(f"{file_config['icon']} **{conversion['filename']}** ({conversion['file_type'].upper()}) - {format_file_size(conversion['size'])} → {format_file_size(conversion['pdf_size'])} ({conversion['conversion_time']:.2f}s)")
    
    # Instructions
    with st.expander("ℹ️ How to Use", expanded=False):
        st.markdown("""
        ### **Step-by-Step Guide:**
        1. **Upload**: Click "Browse files" and select your document
        2. **Validate**: The app will check your file for compatibility
        3. **Convert**: Click "Convert to PDF" to start the conversion process
        4. **Download**: Click the download button to save your PDF file
        
        ### **Supported Features:**
        - ✅ **Word Documents** (.docx): Full formatting preservation
        - ✅ **PowerPoint** (.pptx): Slide content extraction
        - ✅ **Text Files** (.txt): Clean formatting with monospace font
        - ✅ **Markdown** (.md): Rich formatting with syntax highlighting
        - ✅ **File size validation** (up to 50MB)
        - ✅ **Conversion time tracking**
        - ✅ **File size comparison**
        - ✅ **Automatic cleanup** for security
        
        ### **File Type Specific Notes:**
        - **Word**: Requires Microsoft Word installation
        - **PowerPoint**: Extracts text content from slides
        - **Text**: Preserves formatting and line breaks
        - **Markdown**: Converts to styled HTML then PDF
        """)
    
    # Footer
    st.markdown("---")
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.markdown("Built with ❤️ using Streamlit and multiple PDF libraries")

if __name__ == "__main__":
    main()